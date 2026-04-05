from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "LiRx Prescription Transfer System"
    subtitle.text = "Automated Data Logging & Lead Capture\nProject Overview for Stakeholders"

    # --- Slide 2: Home Page Lead Capture ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Step 1: Homepage Lead Capture"
    content = slide.placeholders[1]
    content.text = "• Quick Transfer entry point on the homepage\n• Captures leads instantly with minimal friction\n• Mobile-first, streamlined Tailwind CSS design"
    
    img_path = '/Users/sriramreddy/.gemini/antigravity/brain/ebc96b10-003c-4939-a728-e952174be706/presentation_home_form_1775360448214.png'
    slide.shapes.add_picture(img_path, Inches(5), Inches(2), height=Inches(4))

    # --- Slide 3: Detailed Data Intake ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Step 2: Comprehensive Transfer Form"
    content = slide.placeholders[1]
    content.text = "• In-depth medical and pharmacy data collection\n• New Rx Number and Medication list fields\n• Secure data verification before submission"
    
    img_path = '/Users/sriramreddy/.gemini/antigravity/brain/ebc96b10-003c-4939-a728-e952174be706/presentation_transfer_form_1775360459111.png'
    slide.shapes.add_picture(img_path, Inches(5), Inches(1.5), height=Inches(5))

    # --- Slide 4: Reliable Success State ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Step 3: Verified Success Feedback"
    content = slide.placeholders[1]
    content.text = "• Solved browser CORS/Opaque response blocks\n• Implemented Next.js Server-Side Proxy\n• Professional success UI for patient confidence"
    
    img_path = '/Users/sriramreddy/.gemini/antigravity/brain/ebc96b10-003c-4939-a728-e952174be706/presentation_success_msg_final_capture_1775360850564.png'
    slide.shapes.add_picture(img_path, Inches(5), Inches(2), height=Inches(4))

    # --- Slide 5: Automated Operations Result ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Result: Automated Logging"
    content = slide.placeholders[1]
    content.text = "• Instant sync to central Google Sheet\n• ZERO manual transcription required\n• Columns: Name, Phone, Rx#, Pharmacy, Meds"
    
    img_path = '/Users/sriramreddy/.gemini/antigravity/brain/ebc96b10-003c-4939-a728-e952174be706/presentation_google_sheet_1775360664566.png'
    slide.shapes.add_picture(img_path, Inches(1), Inches(4.5), width=Inches(8))

    # --- Slide 6: Secure Architecture ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Architecture"
    content = slide.placeholders[1]
    content.text = "1. Frontend: Secure Data Intake\n2. Next.js Proxy: Encrypted Data Forwarding\n3. Apps Script: Automated Spreadsheet Sync"
    
    prs.save('/Users/sriramreddy/Desktop/LiRx_Prescription_Transfer_Automation.pptx')
    print("Presentation saved successfully at /Users/sriramreddy/Desktop/LiRx_Prescription_Transfer_Automation.pptx")

if __name__ == "__main__":
    create_presentation()
